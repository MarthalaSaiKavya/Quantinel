"""Build the Quantinel one-pager into an editable PPTX deck (one slide per section)."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (from quantinel_one_pager.html) ----
YC        = RGBColor(0xFF, 0x66, 0x00)
YC_DARK   = RGBColor(0xCC, 0x52, 0x00)
YC_LIGHT  = RGBColor(0xFF, 0xF4, 0xEB)
YC_MID    = RGBColor(0xFF, 0xE0, 0xC7)
BLACK     = RGBColor(0x11, 0x11, 0x11)
G900      = RGBColor(0x1A, 0x1A, 0x1A)
G700      = RGBColor(0x40, 0x40, 0x40)
G600      = RGBColor(0x52, 0x52, 0x52)
G500      = RGBColor(0x73, 0x73, 0x73)
G300      = RGBColor(0xD4, 0xD4, 0xD4)
G200      = RGBColor(0xE5, 0xE5, 0xE5)
G100      = RGBColor(0xF5, 0xF5, 0xF5)
G50       = RGBColor(0xFA, 0xFA, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Space Grotesk"
BODY_FONT  = "Inter"

EMU = 914400
def IN(v): return Emu(int(v * EMU))

prs = Presentation()
prs.slide_width  = IN(13.333)
prs.slide_height = IN(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None, shadow_off=True):
    sp = s.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if shadow_off:
        sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        spacing=1.08, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font, tracking)."""
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(para, dict):
            if para.get("space_before") is not None:
                p.space_before = Pt(para["space_before"])
            if para.get("space_after") is not None:
                p.space_after = Pt(para["space_after"])
            parts = para["parts"]
        else:
            parts = para
        for (t, size, color, bold, font, *rest) in parts:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
            if rest and rest[0]:
                _track(r, rest[0])
    return tb


def _track(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def sec_header(s, number, title, x=0.92, y=0.62):
    txt(s, x, y, 1.0, 0.4, [[(number, 13, YC, True, TITLE_FONT, 0.6)]])
    txt(s, x + 0.62, y - 0.08, 9.0, 0.6, [[(title, 26, BLACK, True, TITLE_FONT)]])
    rect(s, x, y + 0.52, SW - 2 * x, 0.014, fill=G200)


# ============================================================
# SLIDE 1 — TITLE / HEADER
# ============================================================
s = slide(); bg(s, WHITE)
rect(s, 0, 0, SW, 0.22, fill=YC)
rect(s, 0, SH - 0.22, SW, 0.22, fill=BLACK)
# logo
txt(s, 0.92, 1.7, 11.5, 1.3,
    [[("Quantinel", 60, BLACK, True, TITLE_FONT), (".", 60, YC, True, TITLE_FONT)]])
# tagline
txt(s, 0.92, 3.2, 11.4, 2.2, [
    [("Institutional Agentic Quantum Trading Decision Engine", 24, BLACK, True, BODY_FONT)],
    {"space_before": 10, "parts": [
        ("A modular backtesting framework benchmarking ", 17, G600, False, BODY_FONT),
        ("classical Markowitz", 17, YC_DARK, True, BODY_FONT),
        (" against ", 17, G600, False, BODY_FONT),
        ("quantum-augmented forecasting & optimization", 17, YC_DARK, True, BODY_FONT),
        (", with an AI master agent delivering the final auditable verdict.", 17, G600, False, BODY_FONT)]},
], spacing=1.3)
# chips
chips = ["Baseline vs Challenger", "NVDA · GOOG", "xpyq", "Exa", "OpenRouter"]
cx = 0.92
for c in chips:
    w = 0.34 + len(c) * 0.092
    rect(s, cx, 5.65, w, 0.42, fill=YC_LIGHT, line=YC_MID, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, cx, 5.65, w, 0.42, [[(c, 11.5, YC_DARK, True, BODY_FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + 0.18


# ============================================================
# SLIDE 2 — PROBLEM & SOLUTION
# ============================================================
s = slide(); bg(s, WHITE)
sec_header(s, "01", "The Problem & Our Solution")

problems = [
    ("Classical Optimization Hits Walls",
     "Markowitz mean-variance is convex but brittle — covariance estimates degrade under regime shifts and fat tails."),
    ("Quantum Claims Are Unverified",
     "Quantum trading strategies are hyped without rigorous, fair comparisons on identical data and risk pipelines."),
    ("No Unified Decision Layer",
     "Trading desks lack a single explainable output fusing quant scores, market news, and recommendations."),
]
solutions = [
    ("Fair Dual-Engine Benchmark",
     "Baseline & challenger pipelines share the same data, risk, executor, and scorer — isolating forecast & optimizer."),
    ("Remote Quantum Compute via xpyq",
     "SVD factor extraction and QUBO-style optimization offloaded to xpyq with automatic fallback on timeout."),
    ("AI-Powered Final Verdict",
     "MasterAgent fuses Exa real-time sentiment with quant scores into an auditable, plain-English recommendation."),
]

col_w = 5.55
gap = 0.35
x0 = 0.92
y0 = 1.55
card_h = 1.55

def ps_column(s, x, label, items, accent_fill, accent_line, label_color, icon_fill):
    txt(s, x, y0, col_w, 0.34, [[(label, 12, label_color, True, BODY_FONT, 1.2)]])
    cy = y0 + 0.5
    for title, body in items:
        rect(s, x, cy, col_w, card_h, fill=accent_fill, line=accent_line, line_w=1.2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, x + 0.22, cy + 0.24, 0.18, card_h - 0.48, fill=label_color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        txt(s, x + 0.6, cy + 0.22, col_w - 0.85, card_h - 0.4, [
            [(title, 15, BLACK, True, BODY_FONT)],
            {"space_before": 5, "parts": [(body, 12, G600, False, BODY_FONT)]},
        ], spacing=1.16)
        cy += card_h + 0.22

ps_column(s, x0, "PROBLEM", problems, G50, G200, G700, G200)
ps_column(s, x0 + col_w + gap, "SOLUTION", solutions, YC_LIGHT, YC_MID, YC_DARK, YC_MID)


# ============================================================
# SLIDE 3 — SYSTEM ARCHITECTURE
# ============================================================
s = slide(); bg(s, WHITE)
sec_header(s, "02", "System Architecture")

# top data node
def node(s, x, y, w, h, title, sub, fill, line, tcolor, scolor, tsize=14, ssize=10.5):
    rect(s, x, y, w, h, fill=fill, line=line, line_w=1.3,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    txt(s, x + 0.08, y + 0.1, w - 0.16, h - 0.2, [
        [(title, tsize, tcolor, True, BODY_FONT)],
        {"space_before": 2, "parts": [(sub, ssize, scolor, False, BODY_FONT)]},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)

LX = 0.92
RX = SW - 0.92
FULL_W = RX - LX  # ~11.49

def hrow(s, y, steps, descs, box_fill, box_line, name_col, arrow_col, h=0.92):
    """One pipeline rendered as a left-to-right row of step boxes joined by arrows."""
    n = len(steps)
    arrow = 0.3
    bw = (FULL_W - (n - 1) * arrow) / n
    x = LX
    for i, (st, ds) in enumerate(zip(steps, descs)):
        rect(s, x, y, bw, h, fill=box_fill, line=box_line, line_w=1.2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13)
        txt(s, x + 0.06, y + 0.08, bw - 0.12, h - 0.16, [
            [(st, 11, name_col, True, BODY_FONT)],
            {"space_before": 3, "parts": [(ds, 9, G500, False, BODY_FONT)]},
        ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.02)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, IN(x + bw + 0.05),
                                    IN(y + h / 2 - 0.08), IN(arrow - 0.1), IN(0.16))
            ar.fill.solid(); ar.fill.fore_color.rgb = arrow_col; ar.line.fill.background()
            ar.shadow.inherit = False
        x += bw + arrow

def down_arrow(s, cx, y, h=0.2, color=G300):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, IN(cx - 0.1), IN(y), IN(0.2), IN(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    a.shadow.inherit = False

def strip(s, y, fill, label_runs, h=0.36):
    rect(s, LX, y, FULL_W, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, LX + 0.24, y, FULL_W - 0.48, h, [label_runs], anchor=MSO_ANCHOR.MIDDLE)

# --- Band A: data source ---
dw = 3.5
node(s, SW / 2 - dw / 2, 1.46, dw, 0.66, "Portfolio Data", "NVDA / GOOG OHLCV",
     BLACK, BLACK, WHITE, G300, tsize=15)
down_arrow(s, SW / 2, 2.2, h=0.18, color=YC)
txt(s, SW / 2 - 2.5, 2.42, 5.0, 0.26,
    [[("RUN BOTH PIPELINES IN PARALLEL", 10, YC_DARK, True, BODY_FONT, 1.0)]],
    align=PP_ALIGN.CENTER)

# --- Band B: the two parallel pipelines, each a horizontal row ---
n_steps = ["MomentumForecaster", "SampleCovRisk", "MeanVarianceOptimizer", "PaperExecutor", "BacktestScorer"]
n_descs = ["Recent return signal", "Cov · VaR · CVaR", "Markowitz sizing", "Simulated execution", "Return · Sharpe · IC"]
q_steps = ["QuantumForecaster", "SampleCovRisk", "QAOA Optimizer", "PaperExecutor", "BacktestScorer"]
q_descs = ["xpyq SVD factor", "Cov · VaR · CVaR", "xpyq eig / QUBO", "Simulated execution", "Return · Sharpe · IC"]

strip(s, 2.82, G100,
      [("NORMAL PIPELINE", 10, G900, True, BODY_FONT, 1.0),
       ("     classical baseline", 10, G500, False, BODY_FONT)])
hrow(s, 3.26, n_steps, n_descs, WHITE, G300, G900, G300)

strip(s, 4.40, YC_LIGHT,
      [("QUANTUM PIPELINE", 10, YC_DARK, True, BODY_FONT, 1.0),
       ("     xpyq challenger", 10, YC_DARK, False, BODY_FONT)])
hrow(s, 4.84, q_steps, q_descs, YC_LIGHT, YC_MID, YC_DARK, YC_MID)

down_arrow(s, SW / 2, 5.92, h=0.18, color=YC)
txt(s, SW / 2 - 3.0, 6.14, 6.0, 0.24,
    [[("BOTH BRANCHES MERGE INTO THE DECISION LAYER", 9.5, G500, True, BODY_FONT, 1.0)]],
    align=PP_ALIGN.CENTER)

# --- Band C: decision / fusion chain ---
chain = [
    ("Comparison Summary", "Return · Sharpe · accuracy gap", G50, G200, BLACK, G600),
    ("Exa Headlines + Sentiment", "Real-time market intel", G50, G200, BLACK, G600),
    ("OpenRouter MasterAgent", "LLM reasoning engine", YC, YC, WHITE, YC_LIGHT),
    ("Plain-English Decision", "Winner · Rationale · Trace", BLACK, BLACK, WHITE, G300),
]
cn = len(chain)
carrow = 0.32
ccw = (FULL_W - (cn - 1) * carrow) / cn
cy = 6.48
cx = LX
for i, (t, sub, fill, line, tcol, scol) in enumerate(chain):
    node(s, cx, cy, ccw, 0.86, t, sub, fill, line, tcol, scol, tsize=12, ssize=9)
    if i < cn - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, IN(cx + ccw + 0.05),
                                IN(cy + 0.43 - 0.08), IN(carrow - 0.1), IN(0.16))
        ar.fill.solid(); ar.fill.fore_color.rgb = YC; ar.line.fill.background()
        ar.shadow.inherit = False
    cx += ccw + carrow


# ============================================================
# SLIDE 4 — KEY CAPABILITIES
# ============================================================
s = slide(); bg(s, WHITE)
sec_header(s, "03", "Key Capabilities")

feats = [
    ("Quantum SVD Extraction", "QuantumForecaster offloads SVD decomposition to xpyq — isolating the strongest hidden market factor for alpha.", "xpyq Remote Compute", True),
    ("Fair Apples-to-Apples", "Risk, execution, and scoring layers are identical across both pipelines. Only forecaster and optimizer differ.", "Controlled Experiment", False),
    ("Automatic Fallback", "If xpyq times out, the quantum engines fall back to Momentum and DiscreteQUBO — always producing a result.", "Resilient Design", True),
    ("Market Intelligence", "Exa scrapes live headlines and extracts sentiment themes, enriching the MasterAgent's comparison.", "Exa API", True),
    ("Auditable Decision Trace", "Every recommendation outputs return gaps, Sharpe gaps, accuracy, risk breach diffs, and xpyq counts.", "Full Transparency", False),
    ("LLM Master Agent", "MasterAgent fuses sentiment, risk metrics, and comparison into a winner with a plain-English rationale.", "OpenRouter LLM", True),
]

fcw = 3.74
fch = 2.12
fgx = 0.31
fgy = 0.3
fx0 = 0.92
fy0 = 1.6
for i, (title, body, tag, orange) in enumerate(feats):
    col = i % 3; row = i // 3
    x = fx0 + col * (fcw + fgx)
    y = fy0 + row * (fch + fgy)
    rect(s, x, y, fcw, fch, fill=WHITE, line=G200, line_w=1.2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    rect(s, x, y, fcw, 0.09, fill=(YC if orange else G300),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.3, y + 0.28, fcw - 0.6, 0.5, [[(title, 15.5, BLACK, True, BODY_FONT)]])
    txt(s, x + 0.3, y + 0.78, fcw - 0.6, 1.0, [[(body, 11.5, G600, False, BODY_FONT)]], spacing=1.18)
    tagw = 0.3 + len(tag) * 0.083
    rect(s, x + 0.3, y + fch - 0.5, tagw, 0.32,
         fill=(YC_LIGHT if orange else G100), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.3, y + fch - 0.5, tagw, 0.32,
        [[(tag, 9.5, (YC_DARK if orange else G600), True, BODY_FONT, 0.5)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 5 — ROADMAP
# ============================================================
s = slide(); bg(s, WHITE)
sec_header(s, "04", "Roadmap")

phases = [
    ("NOW", "Foundation", YC, WHITE,
     ["NVDA / GOOG OHLCV data", "Dual-engine backtester", "xpyq SVD + QUBO optimization",
      "Exa + OpenRouter AI agent", "Weekly rebalance loop"]),
    ("NEXT", "Expansion", YC_MID, YC_DARK,
     ["Live market data feeds", "Multi-asset universe (S&P 500)", "Deeper quantum circuits",
      "Risk-adjusted alpha tracking", "Portfolio constraint tuning"]),
    ("VISION", "Scale", G200, G700,
     ["Real hardware QPU execution", "Institutional-grade backtests", "Multi-agent hedge fund layer",
      "Explainable AI audit trail", "Open research publication"]),
]
pcw = 3.74
pgx = 0.31
px0 = 0.92
py0 = 1.7
pch = 3.7
for i, (badge, title, bfill, btext, items) in enumerate(phases):
    x = px0 + i * (pcw + pgx)
    rect(s, x, py0, pcw, pch, fill=G50, line=G200, line_w=1.2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    bw = 0.5 + len(badge) * 0.1
    rect(s, x + 0.32, py0 + 0.32, bw, 0.4, fill=bfill,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.32, py0 + 0.32, bw, 0.4, [[(badge, 11, btext, True, BODY_FONT, 0.8)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.32, py0 + 0.9, pcw - 0.6, 0.5, [[(title, 20, BLACK, True, TITLE_FONT)]])
    items_runs = []
    for it in items:
        items_runs.append({"space_before": 7, "parts": [
            ("—  ", 12, YC, True, BODY_FONT), (it, 12.5, G700, False, BODY_FONT)]})
    txt(s, x + 0.32, py0 + 1.5, pcw - 0.6, pch - 1.7, items_runs, spacing=1.1)


# ============================================================
# SLIDE 6 — CLOSING CTA
# ============================================================
s = slide(); bg(s, BLACK)
rect(s, 0, 0, SW, 0.22, fill=YC)
txt(s, 1.2, 2.5, 10.9, 1.6, [
    [("Agentic quantum risk decisions.", 40, WHITE, True, TITLE_FONT)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
txt(s, 1.8, 3.95, 9.7, 1.2, [
    [("Baseline and challenger engines on the same market, risk models, and live news. ",
      17, G300, False, BODY_FONT),
     ("Every number in the open. Every decision auditable.", 17, YC, True, BODY_FONT)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.35)
txt(s, 1.2, 5.7, 10.9, 0.5,
    [[("Quantinel", 18, WHITE, True, TITLE_FONT), (".", 18, YC, True, TITLE_FONT)]],
    align=PP_ALIGN.CENTER)
txt(s, 1.2, 6.15, 10.9, 0.4,
    [[("github.com/MarthalaSaiKavya/Quantinel", 12, G500, False, BODY_FONT)]],
    align=PP_ALIGN.CENTER)

out = "/Users/kalidindiadithya/Desktop/Quantinel/docs/Quantinel_OnePager.pptx"
prs.save(out)
print("Saved", out, "slides:", len(prs.slides._sldIdLst))
